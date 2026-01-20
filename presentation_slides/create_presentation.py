"""
Generate PowerPoint presentation for DoED Public Comment Analysis Pipeline
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_title_slide(prs, title, subtitle):
    """Create title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    return slide

def create_content_slide(prs, title, content_items):
    """Create bullet point slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    
    text_frame = slide.placeholders[1].text_frame
    text_frame.clear()
    
    for item in content_items:
        p = text_frame.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(18)
    
    return slide

def create_two_column_slide(prs, title, left_items, right_items):
    """Create two-column slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    
    # Remove default placeholder
    for shape in slide.shapes:
        if shape.placeholder_format.idx == 1:
            sp = shape.element
            sp.getparent().remove(sp)
    
    # Add left column
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(5))
    text_frame = left_box.text_frame
    for item in left_items:
        p = text_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
    
    # Add right column
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.5), Inches(4.5), Inches(5))
    text_frame = right_box.text_frame
    for item in right_items:
        p = text_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
    
    return slide

def main():
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title
    create_title_slide(
        prs,
        "DoED Public Comment Analysis Pipeline",
        "Automated Analysis of Regulations.gov Comments\nUsing Azure AI Services"
    )
    
    # Slide 2: The Challenge
    create_content_slide(
        prs,
        "The Challenge",
        [
            "95+ public comments received on IDEA regulation changes",
            "Most comments include PDF/DOCX attachments (88 of 95)",
            "Manual review is time-consuming and inconsistent",
            "Need to identify themes, sentiment, and key recommendations",
            "Requirement for comprehensive, unbiased analysis"
        ]
    )
    
    # Slide 3: Solution Overview
    create_content_slide(
        prs,
        "Solution Overview",
        [
            "Automated 3-phase pipeline leveraging Azure AI",
            "Phase 1: Fetch comments from Regulations.gov API",
            "Phase 2: Extract text from PDF/DOCX attachments",
            "Phase 3: AI-powered categorization and analysis",
            "Outputs: Structured categorizations + collective analysis"
        ]
    )
    
    # Slide 4: Architecture
    create_content_slide(
        prs,
        "Solution Architecture",
        [
            "📥 Regulations.gov API → Comment Collection",
            "📄 Azure Document Intelligence → Text Extraction (OCR)",
            "🤖 Azure AI Agents → Categorization & Analysis",
            "📊 Structured Output → JSON categorizations + CSV data",
            "🔄 Automated workflow with organized output folders"
        ]
    )
    
    # Slide 5: Phase 1 - Data Collection
    create_content_slide(
        prs,
        "Phase 1: Data Collection",
        [
            "Script: fetch_regulations_comments.py",
            "Connects to Regulations.gov API with API key",
            "Fetches all comments for specified docket",
            "Extracts metadata: commenter, date, organization",
            "Identifies comments with attachments",
            "Output: JSON files in output/1_fetch/"
        ]
    )
    
    # Slide 6: Phase 2 - Text Extraction
    create_content_slide(
        prs,
        "Phase 2: Text Extraction & Consolidation",
        [
            "Script: consolidate_comments_to_csv.py",
            "Downloads PDF/DOCX attachments from regulations.gov",
            "Extracts text using PyPDF2 and python-docx",
            "Browser-like headers bypass download restrictions",
            "Combines inline text + attachment text",
            "Output: Consolidated CSV in output/2_consolidate/"
        ]
    )
    
    # Slide 7: Phase 3 - AI Analysis (Part 1)
    create_content_slide(
        prs,
        "Phase 3: AI Analysis - Individual Categorization",
        [
            "Script: process_csv_rows.py",
            "Azure AI Agent categorizes each comment individually",
            "Extracts: sentiment, themes, key arguments",
            "Identifies: commenter type, concerns, recommendations",
            "Streaming responses for real-time progress",
            "Output: categorizations_[timestamp].json"
        ]
    )
    
    # Slide 8: Phase 3 - AI Analysis (Part 2)
    create_content_slide(
        prs,
        "Phase 3: AI Analysis - Collective Analysis",
        [
            "Agent receives categorizations in small batches",
            "Maintains thread context across all batches",
            "Identifies common themes and patterns",
            "Groups similar comments together",
            "Generates final collective analysis report",
            "Output: grouped_analysis_[timestamp].json"
        ]
    )
    
    # Slide 9: Key Features
    create_content_slide(
        prs,
        "Key Features & Benefits",
        [
            "✅ Handles 95+ comments including attachment-only submissions",
            "✅ Automated text extraction from PDFs and Word docs",
            "✅ No manual intervention required after initial setup",
            "✅ Organized output structure (output/1_fetch, 2_consolidate, 3_analysis)",
            "✅ Automatic file linking between pipeline phases",
            "✅ Configurable batch sizes and row limits for testing"
        ]
    )
    
    # Slide 10: Technology Stack
    create_two_column_slide(
        prs,
        "Technology Stack",
        [
            "Azure Services:",
            "• Azure AI Agents",
            "• Azure OpenAI",
            "• Azure CLI Authentication",
            "",
            "Python Libraries:",
            "• Semantic Kernel",
            "• PyPDF2",
            "• python-docx",
            "• requests"
        ],
        [
            "APIs:",
            "• Regulations.gov API",
            "",
            "Data Formats:",
            "• JSON for categorizations",
            "• CSV for consolidated text",
            "",
            "Output Structure:",
            "• output/1_fetch/",
            "• output/2_consolidate/",
            "• output/3_analysis/"
        ]
    )
    
    # Slide 11: Results & Metrics
    create_content_slide(
        prs,
        "Results & Metrics",
        [
            "📊 Successfully processed 95 public comments",
            "📄 Extracted text from 40+ PDF/DOCX attachments",
            "🤖 Generated individual categorizations for all comments",
            "📈 Produced comprehensive collective analysis",
            "⏱️ Reduced analysis time from weeks to hours",
            "💰 Configurable batch processing for cost optimization"
        ]
    )
    
    # Slide 12: Sample Insights
    create_content_slide(
        prs,
        "Sample Analysis Output",
        [
            "Common Themes Identified:",
            "• Opposition to removing disproportionality data collection",
            "• Concerns about civil rights protections",
            "• Need for transparency and accountability",
            "• Historical discrimination concerns",
            "",
            "Commenter Types: Educators, advocacy organizations, parents",
            "Sentiment: Predominantly opposed to proposed changes"
        ]
    )
    
    # Slide 13: Workflow Automation
    create_content_slide(
        prs,
        "Fully Automated Workflow",
        [
            "1. Run: python fetch_regulations_comments.py",
            "   → Fetches comments, saves to output/1_fetch/",
            "",
            "2. Run: python consolidate_comments_to_csv.py",
            "   → Auto-finds latest JSON, extracts attachments",
            "   → Saves CSV to output/2_consolidate/",
            "",
            "3. Run: python process_csv_rows.py",
            "   → Auto-finds latest CSV, runs AI analysis",
            "   → Saves results to output/3_analysis/"
        ]
    )
    
    # Slide 14: Configuration Options
    create_content_slide(
        prs,
        "Configuration & Customization",
        [
            "Environment Variables (.env):",
            "• REGULATIONS_GOV_API_KEY",
            "• Azure credentials (via Azure CLI)",
            "",
            "Script Parameters:",
            "• max_rows: Limit comments for testing (None = all)",
            "• batch_size: Control AI agent batch processing (2-20)",
            "",
            "Flexible for different use cases and scale"
        ]
    )
    
    # Slide 15: Next Steps & Enhancements
    create_content_slide(
        prs,
        "Potential Enhancements",
        [
            "🔄 Add summarization for long comments (Azure AI Language)",
            "📊 Create visualization dashboard for results",
            "🔍 Implement semantic search across comments",
            "📧 Email notifications on completion",
            "🌐 Web interface for non-technical users",
            "🔐 Enhanced authentication and security"
        ]
    )
    
    # Slide 16: Cost Optimization
    create_content_slide(
        prs,
        "Cost Optimization Strategies",
        [
            "Batch Processing: Smaller batches = more calls but safer",
            "Row Limiting: Test with max_rows before full run",
            "Caching: Reuse Phase 1 & 2 outputs for multiple analyses",
            "Thread Reuse: Single thread maintains context across batches",
            "Incremental Processing: Process new comments only",
            "Configurable agent IDs: Use different models for cost/quality balance"
        ]
    )
    
    # Slide 17: Lessons Learned
    create_content_slide(
        prs,
        "Key Lessons Learned",
        [
            "API Challenges: Regulations.gov blocks direct downloads",
            "Solution: Browser-like headers bypass restrictions",
            "",
            "Token Limits: Can't send all comments at once",
            "Solution: Batch processing with thread context",
            "",
            "File Organization: Output files scattered initially",
            "Solution: Organized folder structure (output/1_fetch, etc.)"
        ]
    )
    
    # Slide 18: Thank You
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Add title
    left = Inches(1)
    top = Inches(2.5)
    width = Inches(8)
    height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Thank You"
    p.font.size = Pt(54)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    # Add subtitle
    left = Inches(1)
    top = Inches(4)
    width = Inches(8)
    height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Questions?"
    p.font.size = Pt(36)
    p.alignment = PP_ALIGN.CENTER
    
    # Save presentation
    output_file = "DoED_Comment_Analysis_Presentation.pptx"
    prs.save(output_file)
    print(f"\n✅ Presentation created successfully: {output_file}")
    print(f"📊 Total slides: {len(prs.slides)}")

if __name__ == "__main__":
    main()
